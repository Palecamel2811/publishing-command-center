"""
Document ingestion service.

Handles reading, parsing, chunking, and storing diverse document formats:
- PDF (split sheets, royalty statements, contracts)
- CSV/Excel (DSP reports, royalty exports)
- Text files
- PowerPoint (presentation-based split sheets)
"""

from __future__ import annotations

import csv
import io
import json
import logging
import re
import uuid
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Optional

from openai import OpenAI

from ..config import Settings
from ..models import (
    DocumentType,
    IngestResult,
    RoyaltyEntry,
    Split,
    Work,
)
from ..rag.chunker import LegalFinancialChunker
from ..rag.embedder import EmbeddingService
from ..rag.store import VectorStoreManager
from ..db.database import Session, engine
from ..db.models import Work, Split, RoyaltyEntry as RelRoyaltyEntry, SyncLicense as RelSyncLicense, DocumentChunk
from sqlmodel import select
from sqlalchemy.orm import selectinload

logger = logging.getLogger(__name__)


@dataclass
class IngestionStats:
    """Statistics for a document ingestion."""
    document_id: str
    filename: str
    doc_type: str = "unknown"
    pages: int = 0
    chunks_created: int = 0
    works_found: list[str] = field(default_factory=list)
    splits_found: list[dict] = field(default_factory=list)
    royalties_found: list[dict] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)
    processing_time_ms: int = 0


class DocumentIngestionService:
    """
    Service for ingesting and processing publishing documents.
    
    Workflow:
    1. Detect document type from content/metadata
    2. Extract structured data (works, splits, royalties)
    3. Chunk documents with domain-aware strategies
    4. Generate embeddings
    5. Store in vector database with metadata
    6. Return structured results for database indexing
    """

    # File extension to document type mapping
    EXTENSION_MAP = {
        ".pdf": "pdf",
        ".csv": "csv",
        ".xlsx": "xlsx",
        ".xls": "xls",
        ".txt": "txt",
        ".docx": "docx",
        ".pptx": "pptx",
    }

    # Content-based type detection patterns
    SPLIT_SHEET_PATTERNS = [
        r"split\s*sheet",
        r"writing\s+share",
        r"publishing\s+share",
        r"master\s+rights\s+share",
        r"publisher\s+share\s*:\s*\d+%",
        r"pro:\s*(ASCAP|BMI|SESAC|GEMA|PRS|SOCAN)",
        r"(?:writer|publisher|producer)\s*%\s*(?:share)",
    ]
    
    ROYALTY_STATEMENT_PATTERNS = [
        r"royalty\s*(statement|report|reporting)",
        r"mechanical\s*royalty",
        r"performance\s*royalty",
        r"streaming\s+revenue",
        r"unit\s+(sale|stream)",
        r"gross\s*(earnings|revenue|amount)",
        r"net\s*(earnings|revenue|amount)",
    ]
    
    CONTRACT_PATTERNS = [
        r"(?:licensing|agreement|contract)",
        r"term\s*of\s*this\s+agreement",
        r"rights\s+granted",
        r"royalty\s+rate",
        r"territory",
        r"herein\s+agreed",
    ]

    def __init__(
        self,
        settings: Settings,
        embedder: EmbeddingService,
        store: VectorStoreManager,
        llm_client: Optional[OpenAI] = None,
    ):
        self.settings = settings
        self.embedder = embedder
        self.store = store
        self.llm_client = llm_client
        self.chunker = LegalFinancialChunker(
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap,
        )

    def _clean_work_title(self, title: str | None) -> str | None:
        if not title:
            return None
        title_clean = title.strip().strip('"').strip('“').strip('”')
        # If it starts with & or contains metadata structural tags
        if title_clean.startswith("&"):
            m = re.search(r"Single Title\s*[:–-]?\s*(.*?)\s*(?:Album Title|Primary Artist|ISRC|ISWC|$)", title_clean, re.IGNORECASE)
            if m and len(m.group(1).strip()) > 1:
                return m.group(1).strip()
            return None
        # Strip common noise suffixes if found from filenames or headings
        for noise in ["Split Sheet Report", "Split Sheet", "Royalty Statement", "Sync Contract", "Statement Report", "Report", "Statement", "Splits", "Master Splits"]:
            pattern = re.compile(rf"[ \t_-]*\b{re.escape(noise)}\b[ \t_-]*.*", re.IGNORECASE)
            title_clean = pattern.sub("", title_clean).strip()
        title_clean = re.sub(r"[ \t_-]+\b(LUV|L\.U\.V\.|v\d+)\b.*$", "", title_clean, flags=re.IGNORECASE).strip()
        if len(title_clean) >= 2:
            return title_clean.title() if title_clean.isupper() or "_" in title_clean else title_clean
        return None

    async def ingest_file(
        self,
        file_path: str,
        file_name: str,
        file_data: bytes,
        doc_type: Optional[str] = None,
        work_id: Optional[str] = None,
        db_session: Optional[Session] = None,
    ) -> IngestResult:
        """
        Ingest a single file into the vector store.
        
        Args:
            file_path: Path to save the original file
            file_name: Original filename for display
            file_data: Raw file bytes
            doc_type: Override document type detection
            work_id: Associate with existing work
        
        Returns:
            IngestResult with structured findings
        """
        import time
        start = time.time()
        
        document_id = str(uuid.uuid4())
        result = IngestResult(
            document_id=document_id,
            chunks_created=0,
            works_found=[],
            splits_found=[],
            royalties_found=[],
        )
        
        try:
            # 0. Save original file to disk if file_path is specified
            if file_path:
                try:
                    target_p = Path(file_path)
                    target_p.parent.mkdir(parents=True, exist_ok=True)
                    with open(target_p, "wb") as f:
                        f.write(file_data)
                except Exception as save_err:
                    logger.warning(f"Could not save file to {file_path}: {save_err}")

            # 1. Detect document type
            detected_type = doc_type or self._detect_type(file_name, file_data)
            result_doc_type = self._normalize_doc_type(detected_type)
            
            # 2. Parse file content
            text_content, metadata = await self._parse_file(
                file_name, file_data, result_doc_type
            )
            
            if not text_content:
                result.warnings.append(f"Could not extract text from {file_name}")
                return result
            
            # 3. Chunk the document
            chunks = self.chunker.chunk_document(
                text=text_content,
                doc_type=result_doc_type,
                source_filename=file_name,
                work_id=work_id,
            )
            
            # 4. Extract structured data during chunking
            work_titles = set()
            split_data = []
            royalty_data = []
            
            isrc_match = re.search(r"ISRC\b.*?\b([A-Z]{2}-?[A-Z0-9]{3}-?[0-9]{2}-?[0-9]{5}|[A-Z]{2}[A-Z0-9]{10})", text_content, re.IGNORECASE)
            iswc_match = re.search(r"ISWC\b.*?\b(T-?\d{3}\.?\d{3}\.?\d{3}-?\d|T\d{9}\d)", text_content, re.IGNORECASE)
            isrc_val = isrc_match.group(1).replace("-", "").strip().upper() if isrc_match else None
            iswc_val = iswc_match.group(1).strip().upper() if iswc_match else None

            for chunk_text, chunk_meta in chunks:
                # Extract work titles
                if chunk_meta.work_title and chunk_meta.work_title not in ("Unknown", ""):
                    work_titles.add(chunk_meta.work_title)

            # Extract split data from split sheets table
            if "split" in result_doc_type or "split" in file_name.lower():
                pro_map = {}
                for line in text_content.splitlines():
                    # Parse bulleted PRO entries: e.g. "  • Jordan Lee: BMI"
                    pro_m = re.search(r"•\s*([A-Za-z\s]+?):\s*([A-Z]{3,5})", line)
                    if pro_m:
                        pro_map[pro_m.group(1).strip().lower()] = pro_m.group(2).strip().upper()
                
                has_pipes = False
                for line in text_content.splitlines():
                    if "%" in line and "|" in line:
                        has_pipes = True
                        parts = [p.strip() for p in line.split("|") if p.strip()]
                        if len(parts) >= 2:
                            name = parts[0]
                            if name.lower() not in ("party name", "party", "total share", "total", "parties", "share", "---", "role", "pro"):
                                share_match = re.search(r"(\d+(?:\.\d+)?)", parts[1])
                                if share_match:
                                    share_pct = float(share_match.group(1))
                                    role = parts[2] if len(parts) >= 3 else None
                                    
                                    p_key = name.lower()
                                    pro_val = pro_map.get(p_key)
                                    if not pro_val and "you" in p_key:
                                        pro_val = pro_map.get("you")
                                        
                                    split_data.append({
                                        "party_name": name,
                                        "share_percentage": share_pct,
                                        "notes": role,
                                        "pro": pro_val,
                                        "source_document": file_name,
                                    })
                
                if not has_pipes:
                    # Parse space-separated layout: name, role, writer%, publisher%, PRO, IPI
                    pattern = r'([A-Za-z\s&.\-(): Palestine \"\u201c\u201d]{2,100}?)\s+(Main Writer|Co-Writer|Publisher|Admin|Writer|Producer)\s+(\d+(?:\.\d+)?)%\s+(\d+(?:\.\d+)?)%\s+(ASCAP|BMI|SESAC|PRS|GEMA|SACEM|-)\s+(\d+|N/A)'
                    for match in re.finditer(pattern, text_content, re.IGNORECASE):
                        name = match.group(1).strip().strip('\n').strip()
                        role = match.group(2).strip()
                        writer_pct = float(match.group(3))
                        pub_pct = float(match.group(4))
                        pro = match.group(5).strip()
                        pro_val = pro.upper() if pro != "-" else None
                        
                        # Clean prefix text if any column headers were matched in flow layout
                        name_lower = name.lower()
                        for prefix in ("pro affiliation ipi number", "ipi number", "affiliation ipi number", "pro affiliation"):
                            if name_lower.startswith(prefix):
                                name = name[len(prefix):].strip()
                                name_lower = name.lower()
                                break
                        
                        # Use publisher % (or writer %) for composition splits
                        share_pct = pub_pct if pub_pct > 0 else writer_pct
                        
                        split_data.append({
                            "party_name": name,
                            "share_percentage": share_pct,
                            "notes": role,
                            "pro": pro_val,
                            "source_document": file_name,
                        })

                    # Also check for PRO-anchored split format: e.g. "Lucas Vance Lead Vocals BMI • #0098... 40.00% $937,082"
                    if not split_data:
                        pattern_pro = r"([A-Za-z\s\"“\.\(\)\-]+?)(?:Lead Vocals|Primary|Featured|Main Producer|Co-Producer|Topline Vocal|Songwriter|Producer|Vocals|Arrangement|Mixing|Mastering)[^%]*?(BMI|ASCAP|SESAC|PRS|GEMA|SACEM)\s*[•\-\:]*\s*#?\d+[^%]*?(\d+(?:\.\d+)?)%\s*\$([0-9,]+(?:\.[0-9]+)?)"
                        for m in re.finditer(pattern_pro, text_content, re.IGNORECASE):
                            raw_name = m.group(1).strip()
                            raw_name = re.sub(r"^[\s\)\(\n\/]+", "", raw_name)
                            lines = [l.strip() for l in raw_name.split("\n") if l.strip() and not any(k in l.upper() for k in ["CONTRIBUTOR", "NAME", "ALIAS", "CONFIDENTIAL", "PAGE", "PRO AFFIL", "ROLE"])]
                            p_name = " ".join(lines).strip()
                            if p_name:
                                pro_val = m.group(2).upper()
                                share_pct = float(m.group(3))
                                split_data.append({
                                    "party_name": p_name,
                                    "share_percentage": share_pct,
                                    "notes": "Songwriter Split",
                                    "pro": pro_val,
                                    "source_document": file_name,
                                })

            # If work_title not found from chunks, search text content first, then filename
            if not work_titles:
                title_match = re.search(r"title\s*[:–-]?\s*[\"“]([^\"”\n\r]{2,50})[\"”]", text_content, re.IGNORECASE)
                if not title_match:
                    title_match = re.search(r"[\"“]([^\"”\n\r]{2,50})[\"”]\s*\n\s*(?:Performed by|recorded by|by)", text_content, re.IGNORECASE)
                if not title_match:
                    title_match = re.search(r"Single Title\s*[:–-]?\s*(.*?)\s*(?:Album Title|Primary Artist|ISRC|ISWC)", text_content, re.IGNORECASE)
                if not title_match:
                    title_match = re.search(r"(?:Song Title|Work Title|Composition Title|Composition|Work|Song|Title)[ \t]*[:–-]?[ \t]*([^\n|,:–-]{2,50})", text_content, re.IGNORECASE)
                if title_match:
                    clean_t = self._clean_work_title(title_match.group(1))
                    if clean_t:
                        work_titles.add(clean_t)
                else:
                    derived_title = (
                        Path(file_name).stem
                        .replace(".pdf.txt", "").replace(".txt", "").replace(".csv", "").replace(".pdf", "")
                        .replace("_sync_contract", "").replace("_split", "")
                        .split("_spotify")[0].split("_apple")[0].split("_youtube")[0].split("_q")[0]
                    )
                    clean_d = self._clean_work_title(derived_title)
                    if clean_d:
                        work_titles.add(clean_d)

            # Extract royalty data from document content or CSV
            ext = Path(file_name).suffix.lower()
            if ext == ".csv":
                try:
                    csv_text = file_data.decode("utf-8", errors="replace")
                    csv_reader = csv.DictReader(io.StringIO(csv_text))
                    for row in csv_reader:
                        row_work = list(work_titles)[0] if work_titles else "Unknown Work"
                        platform = row.get("Platform") or row.get("platform") or "unknown"
                        r_type = row.get("Type") or row.get("royalty_type") or "streaming"
                        gross_str = (row.get("Gross Revenue") or row.get("gross_revenue") or "0").replace(",", "").replace("$", "")
                        net_str = (row.get("Net Revenue") or row.get("net_revenue") or "0").replace(",", "").replace("$", "")
                        fee_str = (row.get("Fees") or row.get("fees") or "0").replace(",", "").replace("$", "")
                        date_str = row.get("Date") or row.get("date") or ""
                        
                        gross = float(gross_str) if gross_str else 0.0
                        net = float(net_str) if net_str else 0.0
                        fees = float(fee_str) if fee_str else max(0.0, gross - net)
                        
                        royalty_data.append({
                            "work_title": row_work,
                            "platform": platform.lower(),
                            "royalty_type": r_type.lower(),
                            "period_start": date_str,
                            "period_end": date_str,
                            "gross_amount": gross,
                            "fees_deducted": fees,
                            "net_amount": net,
                            "source_document": file_name,
                        })
                except Exception as csv_err:
                    logger.warning(f"Failed to parse CSV rows from {file_name}: {csv_err}")
            elif result_doc_type in ("royalty_statement", "split_sheet") or "royalty" in file_name.lower() or "statement" in file_name.lower() or "gross" in text_content.lower():
                gross_match = re.search(r"total\s*gross\b[^\n]*?\$([0-9]{1,3}(?:,[0-9]{3})*(?:\.[0-9]+)?|[0-9]+(?:\.[0-9]+)?)", text_content, re.IGNORECASE)
                if not gross_match:
                    gross_match = re.search(r"gross\b[^\n]*?\$([0-9]{1,3}(?:,[0-9]{3})*(?:\.[0-9]+)?|[0-9]+(?:\.[0-9]+)?)", text_content, re.IGNORECASE)
                
                net_match = re.search(r"total\s*net\b[^\n]*?\$([0-9]{1,3}(?:,[0-9]{3})*(?:\.[0-9]+)?|[0-9]+(?:\.[0-9]+)?)", text_content, re.IGNORECASE)
                if not net_match:
                    net_match = re.search(r"net\b[^\n]*?\$([0-9]{1,3}(?:,[0-9]{3})*(?:\.[0-9]+)?|[0-9]+(?:\.[0-9]+)?)", text_content, re.IGNORECASE)
                
                platform_match = re.search(r"\b(Spotify|Apple\s*Music|YouTube|TikTok|Amazon\s*Music|Deezer|ASCAP|BMI)\b", text_content, re.IGNORECASE)
                period_match = re.search(r"(Q[1-4]\s*\d{4}|\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4})", text_content, re.IGNORECASE)
                
                gross = float(gross_match.group(1).replace(",", "")) if gross_match else 0.0
                net = float(net_match.group(1).replace(",", "")) if net_match else 0.0
                platform = platform_match.group(1).lower().replace(" ", "_") if platform_match else ("spotify" if "spotify" in file_name.lower() else "apple_music" if "apple" in file_name.lower() else "youtube" if "youtube" in file_name.lower() else "other")
                period = period_match.group(1) if period_match else ""

                # First check for multi-platform table breakdown (Spotify, Apple Music, TikTok, YouTube, Sync, Radio, etc.)
                multi_platform_found = False
                for block in re.split(r"(?=\b(?:Spotify|Apple Music|TikTok|YouTube|Sync Licensing|Terrestrial Radio)\b)", text_content):
                    m_plat = re.match(r"\s*(Spotify|Apple Music|TikTok|YouTube|Sync Licensing|Terrestrial Radio)", block, re.IGNORECASE)
                    if m_plat:
                        amounts = re.findall(r"\$([0-9,]+(?:\.[0-9]+)?)", block)
                        if len(amounts) >= 3:
                            multi_platform_found = True
                            p_name = m_plat.group(1).lower().replace(" ", "_")
                            p_gross = float(amounts[0].replace(",", ""))
                            p_fee = float(amounts[1].replace(",", ""))
                            p_net = float(amounts[2].replace(",", ""))
                            royalty_data.append({
                                "work_title": list(work_titles)[0] if work_titles else "Unknown",
                                "platform": p_name,
                                "royalty_type": "sync" if "sync" in p_name else "streaming",
                                "period_start": period,
                                "period_end": period,
                                "gross_amount": p_gross,
                                "fees_deducted": p_fee,
                                "net_amount": p_net,
                                "source_document": file_name,
                            })

                if not multi_platform_found:
                    # Check mechanical & performance breakdowns
                    mech_match = re.search(r"Mechanical\b[^\n]*?\$([0-9]{1,3}(?:,[0-9]{3})*(?:\.[0-9]+)?|[0-9]+(?:\.[0-9]+)?)", text_content, re.IGNORECASE)
                    perf_match = re.search(r"Performance\b[^\n]*?\$([0-9]{1,3}(?:,[0-9]{3})*(?:\.[0-9]+)?|[0-9]+(?:\.[0-9]+)?)", text_content, re.IGNORECASE)
                    
                    fee_pct = max(0.0, (gross - net) / gross) if gross > 0 else 0.0
                    has_granular = False
                    
                    if mech_match:
                        mech_val = float(mech_match.group(1).replace(",", ""))
                        if mech_val > 0:
                            has_granular = True
                            mech_fee = mech_val * fee_pct
                            royalty_data.append({
                                "work_title": list(work_titles)[0] if work_titles else "Unknown",
                                "platform": platform,
                                "royalty_type": "mechanical",
                                "period_start": period,
                                "period_end": period,
                                "gross_amount": mech_val,
                                "fees_deducted": mech_fee,
                                "net_amount": mech_val - mech_fee,
                                "source_document": file_name,
                            })
                    
                    if perf_match:
                        perf_val = float(perf_match.group(1).replace(",", ""))
                        if perf_val > 0:
                            has_granular = True
                            perf_fee = perf_val * fee_pct
                            royalty_data.append({
                                "work_title": list(work_titles)[0] if work_titles else "Unknown",
                                "platform": platform,
                                "royalty_type": "performance",
                                "period_start": period,
                                "period_end": period,
                                "gross_amount": perf_val,
                                "fees_deducted": perf_fee,
                                "net_amount": perf_val - perf_fee,
                                "source_document": file_name,
                            })
                    
                    if not has_granular and (gross > 0 or net > 0):
                        royalty_data.append({
                            "work_title": list(work_titles)[0] if work_titles else "Unknown",
                            "platform": platform,
                            "royalty_type": "streaming",
                            "period_start": period,
                            "period_end": period,
                            "gross_amount": gross,
                            "fees_deducted": max(0.0, gross - net),
                            "net_amount": net,
                            "source_document": file_name,
                        })
            
            # Extract sync licenses if contract
            sync_data = []
            if result_doc_type in ("contract", "license_contract") or "contract" in file_name.lower() or "license" in file_name.lower():
                licensee_match = re.search(r"licensee:\s*([^\n|,]{2,100})", text_content, re.IGNORECASE)
                media_match = re.search(r"media\s*type:\s*([^\n|,]{2,100})", text_content, re.IGNORECASE)
                territory_match = re.search(r"territory:\s*([^\n|]{2,100})", text_content, re.IGNORECASE)
                term_match = re.search(r"term:\s*([^\n|,]{2,100})", text_content, re.IGNORECASE)
                fee_match = re.search(r"license\s*fee:\s*\$?([\d,]+\.?\d*)", text_content, re.IGNORECASE)
                
                contract_title = list(work_titles)[0] if work_titles else Path(file_name).name.replace(".pdf.txt", "").replace(".txt", "").replace("_sync_contract", "").replace("_", " ").title()
                if contract_title not in work_titles:
                    work_titles.add(contract_title)
                
                sync_data.append({
                    "work_title": contract_title,
                    "licensee": licensee_match.group(1).strip() if licensee_match else "Commercial Licensee",
                    "media_type": media_match.group(1).strip() if media_match else "Commercial Advertisement",
                    "territory": territory_match.group(1).strip() if territory_match else "Worldwide",
                    "term_end": term_match.group(1).strip() if term_match else None,
                    "fee": float(fee_match.group(1).replace(",", "")) if fee_match else 0.0,
                })
            
            # Sanitize and clean work_titles
            sanitized_titles = set()
            for t in work_titles:
                c = self._clean_work_title(t)
                if c and len(c) >= 2:
                    sanitized_titles.add(c)
            work_titles = sanitized_titles

            result.works_found = list(work_titles)
            result.splits_found = split_data
            result.royalties_found = royalty_data
            
            # 4.5. Upsert structured data into relational DB
            if db_session:
                try:
                    for title in work_titles:
                        # Find or create work by title
                        work_stmt = select(Work).where(Work.title == title).limit(1)
                        work = db_session.exec(work_stmt).first()
                        if not work:
                            work = Work(title=title, isrc=isrc_val, iswc=iswc_val)
                            db_session.add(work)
                        else:
                            if isrc_val:
                                work.isrc = isrc_val
                            if iswc_val:
                                work.iswc = iswc_val
                            db_session.add(work)
                        db_session.flush()
                        
                        # Upsert splits (skip if party already exists for this work)
                        existing_splits = set()
                        if work.splits:
                            existing_splits = {s.party_name for s in work.splits}
                        for split in split_data:
                            if split["party_name"] not in existing_splits:
                                s = Split(
                                    work_id=work.id,
                                    party_name=split["party_name"],
                                    share_percentage=split["share_percentage"],
                                    pro=split.get("pro"),
                                    notes=split.get("notes"),
                                )
                                db_session.add(s)
                                existing_splits.add(split["party_name"])
                    
                    # Upsert royalty entries from extracted data (with financial validation & deduplication)
                    for entry in royalty_data:
                        work_stmt = select(Work).where(Work.title == entry.get("work_title", ""))
                        work = db_session.exec(work_stmt).first()
                        if work:
                            gross_val = entry.get("gross_amount", 0.0)
                            fees_val = entry.get("fees_deducted", 0.0)
                            net_val = entry.get("net_amount", 0.0)
                            
                            # Financial Rule: Gross MUST be >= Net
                            if gross_val < net_val:
                                gross_val = net_val + fees_val
                            
                            # Deduplication check: do not insert duplicate royalty entry for same work, platform, period, net_amount
                            existing_royalty = db_session.exec(
                                select(RelRoyaltyEntry).where(
                                    RelRoyaltyEntry.work_id == work.id,
                                    RelRoyaltyEntry.platform == entry.get("platform", "unknown"),
                                    RelRoyaltyEntry.period_start == entry.get("period_start", ""),
                                    RelRoyaltyEntry.net_amount == net_val,
                                )
                            ).first()
                            
                            if not existing_royalty:
                                royalty = RelRoyaltyEntry(
                                    work_id=work.id,
                                    platform=entry.get("platform", "unknown"),
                                    royalty_type=entry.get("royalty_type", "other"),
                                    period_start=entry.get("period_start", ""),
                                    period_end=entry.get("period_end", ""),
                                    gross_amount=gross_val,
                                    fees_deducted=fees_val,
                                    net_amount=net_val,
                                    currency=entry.get("currency", "USD"),
                                    source_document=entry.get("source_document", file_name),
                                )
                                db_session.add(royalty)
                    
                    # Upsert sync licenses
                    for s_item in sync_data:
                        work_stmt = select(Work).where(Work.title == s_item.get("work_title", ""))
                        work = db_session.exec(work_stmt).first()
                        if not work:
                            work = Work(title=s_item.get("work_title", ""))
                            db_session.add(work)
                            db_session.flush()
                        
                        existing_sync = db_session.exec(
                            select(RelSyncLicense).where(
                                RelSyncLicense.work_id == work.id,
                                RelSyncLicense.licensee == s_item.get("licensee", "")
                            )
                        ).first()
                        if not existing_sync:
                            sync = RelSyncLicense(
                                work_id=work.id,
                                title=f"Sync License - {work.title}",
                                licensee=s_item.get("licensee", "Licensee"),
                                media_type=s_item.get("media_type", "commercial"),
                                territory=s_item.get("territory", "Worldwide"),
                                term_end=s_item.get("term_end"),
                                fee=s_item.get("fee", 0.0),
                                currency="USD",
                                status="active",
                            )
                            db_session.add(sync)
                    
                    # Record DocumentChunk for history tracking
                    for i, (chunk_text, chunk_meta) in enumerate(chunks):
                        doc_chunk = DocumentChunk(
                            doc_id=document_id,
                            doc_type=result_doc_type,
                            source_filename=file_name,
                            content=chunk_text[:1000],
                            work_title=chunk_meta.work_title or (list(work_titles)[0] if work_titles else None),
                            chunk_index=i,
                        )
                        db_session.add(doc_chunk)
                    
                    # Update work total_earnings
                    for title in work_titles:
                        work = db_session.exec(select(Work).where(Work.title == title)).first()
                        if work:
                            r_sum = sum(r.net_amount for r in work.royalties) if work.royalties else 0.0
                            s_sum = sum(s.fee * 0.5 for s in work.sync_licenses) if work.sync_licenses else 0.0
                            work.total_earnings = r_sum + s_sum
                            db_session.add(work)
                    
                    db_session.commit()
                except Exception as upsert_err:
                    logger.warning(f"DB upsert failed for {file_name}: {upsert_err}")
                    db_session.rollback()
            
            # 5. Generate embeddings for all chunks
            if not chunks:
                result.warnings.append(f"Chunking produced 0 chunks for {file_name}")
                return result

            chunk_texts = [c[0] for c in chunks]
            embeddings = []
            if self.embedder is not None:
                try:
                    embeddings = self.embedder.embed_batch(chunk_texts)
                except Exception as emb_err:
                    logger.warning(f"Embedding failed for {file_name}, using zero-vector fallback: {emb_err}")
                    result.warnings.append(f"Embedding fallback used: {emb_err}")
                    embeddings = [[0.0] * 768 for _ in chunks]
            else:
                embeddings = [[0.0] * 768 for _ in chunks]
            
            # 6. Prepare documents for vector store
            documents = []
            for i, ((chunk_text, chunk_meta), embedding) in enumerate(zip(chunks, embeddings)):
                # Helper to safely serialize metadata values
                def _serialize(val):
                    if val is None:
                        return None
                    if isinstance(val, (str, int, float, bool)):
                        return val
                    # Serialize lists/dicts to JSON string for ChromaDB
                    try:
                        return json.dumps(val)
                    except TypeError:
                        return str(val)

                # Build metadata
                raw_metadata = {
                    "doc_id": document_id,
                    "doc_type": result_doc_type,
                    "source_filename": file_name,
                    "work_id": work_id or chunk_meta.work_title,
                    "period_start": chunk_meta.period_start,
                    "period_end": chunk_meta.period_end,
                    "platform": chunk_meta.platform,
                    "royalty_type": chunk_meta.royalty_type,
                    "chunk_index": i,
                    "parties": chunk_meta.parties,
                    "created_at": datetime.utcnow().isoformat(),
                }
                
                # Filter None and serialize others
                metadata = {
                    k: _serialize(v) for k, v in raw_metadata.items() if _serialize(v) is not None
                }

                doc = {
                    "id": str(uuid.uuid4()),
                    "content": chunk_text,
                    "embedding": embedding,
                    "metadata": metadata,
                }
                documents.append(doc)
            
            # 7. Store in vector database (delete existing chunks for this file first)
            if self.store is not None and len(documents) > 0:
                try:
                    self.store.delete_by_filename(file_name)
                    self.store.add_documents(documents)
                except Exception as store_err:
                    logger.warning(f"Vector store failed for {file_name}: {store_err}")
                    result.warnings.append(f"Vector store warning: {store_err}")
            
            result.chunks_created = len(chunks)
            
            elapsed = int((time.time() - start) * 1000)
            logger.info(
                f"Ingested {file_name}: {len(chunks)} chunks, "
                f"{len(work_titles)} works, {elapsed}ms"
            )
            
        except Exception as e:
            logger.error(f"Ingestion failed for {file_name}: {e}")
            result.warnings.append(f"Error during ingestion: {str(e)}")
        
        return result

    async def _parse_file(
        self,
        file_name: str,
        file_data: bytes,
        doc_type: str,
    ) -> tuple[str, dict]:
        """
        Parse a file and extract text content.
        
        Returns (text_content, metadata_dict).
        """
        ext = Path(file_name).suffix.lower()
        
        if ext == ".pdf":
            return await self._parse_pdf(file_data, doc_type)
        elif ext in (".csv",):
            return await self._parse_csv(file_data, doc_type)
        elif ext in (".xlsx", ".xls"):
            return await self._parse_excel(file_data, doc_type)
        elif ext in (".txt",):
            return file_data.decode("utf-8", errors="replace"), {}
        elif ext == ".docx":
            return await self._parse_docx(file_data, doc_type)
        elif ext == ".pptx":
            return await self._parse_pptx(file_data, doc_type)
        else:
            # Fallback: try as text
            return file_data.decode("utf-8", errors="replace"), {}

    async def _parse_pdf(self, file_data: bytes, doc_type: str) -> tuple[str, dict]:
        """Parse PDF files with pdfplumber tabular extraction fallback."""
        pages = []
        try:
            import pdfplumber
            from io import BytesIO
            with pdfplumber.open(BytesIO(file_data)) as pdf:
                for p in pdf.pages:
                    tables = p.extract_tables()
                    if tables:
                        for table in tables:
                            for row in table:
                                cleaned_row = [str(cell).strip() if cell else "" for cell in row]
                                if any(cleaned_row):
                                    pages.append(" | ".join(cleaned_row))
                    text = p.extract_text()
                    if text and text.strip():
                        pages.append(text.strip())
                if pages:
                    return "\n\n".join(pages), {"pages": len(pdf.pages)}
        except Exception as e:
            logger.info(f"pdfplumber fallback to pypdf: {e}")

        try:
            from pypdf import PdfReader
            from io import BytesIO
            
            reader = PdfReader(BytesIO(file_data))
            for page in reader.pages:
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(text.strip())
            
            metadata = {"pages": len(reader.pages)}
            return "\n\n".join(pages), metadata
            
        except ImportError:
            raise RuntimeError(
                "pypdf not installed. Run: pip install pypdf"
            )


    async def _parse_csv(self, file_data: bytes, doc_type: str) -> tuple[str, dict]:
        """Parse CSV files into a readable text format."""
        text = file_data.decode("utf-8", errors="replace")
        lines = text.strip().split("\n")
        
        if not lines:
            return "", {}
        
        # Parse as CSV with headers
        try:
            reader = csv.reader(io.StringIO(text))
            rows = list(reader)
            
            if not rows:
                return "", {}
            
            headers = [h.strip() for h in rows[0]]
            formatted_lines = [" | ".join(headers)]  # Header row as delimiter-separated
            
            for row in rows[1:100]:  # Limit to first 100 rows
                formatted_lines.append(" | ".join(
                    f"{cell.strip()}" for cell in row
                ))
            
            metadata = {
                "rows": len(rows) - 1,
                "columns": headers,
            }
            
            return "\n".join(formatted_lines), metadata
            
        except Exception:
            # Fallback: treat as plain text
            return text, {}

    async def _parse_excel(self, file_data: bytes, doc_type: str) -> tuple[str, dict]:
        """Parse Excel files into readable text."""
        try:
            from openpyxl import load_workbook
            from io import BytesIO
            
            workbook = load_workbook(BytesIO(file_data), read_only=True, data_only=True)
            sheets = []
            
            for sheet_name in workbook.sheetnames[:5]:  # Limit sheets
                ws = workbook[sheet_name]
                sheet_text = f"Sheet: {sheet_name}\n"
                
                for row in ws.iter_rows(values_only=True):
                    formatted_row = " | ".join(
                        str(cell) if cell is not None else ""
                        for cell in row
                    )
                    if formatted_row.strip():
                        sheet_text += formatted_row + "\n"
                
                sheets.append(sheet_text)
            
            workbook.close()
            
            return "\n\n".join(sheets), {"sheets": workbook.sheetnames}
            
        except ImportError:
            raise RuntimeError(
                "openpyxl not installed. Run: pip install openpyxl"
            )

    async def _parse_docx(self, file_data: bytes, doc_type: str) -> tuple[str, dict]:
        """Parse Word documents."""
        try:
            from docx import Document
            from io import BytesIO
            
            doc = Document(BytesIO(file_data))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            
            return "\n\n".join(paragraphs), {"paragraphs": len(paragraphs)}
            
        except ImportError:
            raise RuntimeError(
                "python-docx not installed. Run: pip install python-docx"
            )

    async def _parse_pptx(self, file_data: bytes, doc_type: str) -> tuple[str, dict]:
        """Parse PowerPoint presentations."""
        try:
            from pptx import Presentation
            from io import BytesIO
            
            prs = Presentation(BytesIO(file_data))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"--- Slide {slide_num} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
            
            return "\n\n".join(text_parts), {"slides": len(prs.slides)}
            
        except ImportError:
            raise RuntimeError(
                "python-pptx not installed. Run: pip install python-pptx"
            )

    def _detect_type(self, file_name: str, file_data: bytes) -> str:
        """Detect document type from filename and content."""
        name_lower = file_name.lower()
        text_sample = file_data[:4096].decode("utf-8", errors="replace").lower()
        
        # Check filename patterns first
        if "split" in name_lower or "splitsheet" in name_lower:
            return "split_sheet"
        if "contract" in name_lower or "license" in name_lower or "agreement" in name_lower:
            return "contract"
        if "royalty" in name_lower or "spotify" in name_lower or "apple" in name_lower or "youtube" in name_lower or "statement" in name_lower:
            return "royalty_statement"
        
        # Check content patterns
        for pattern in self.SPLIT_SHEET_PATTERNS:
            if re.search(pattern, text_sample, re.IGNORECASE):
                return "split_sheet"
        
        for pattern in self.ROYALTY_STATEMENT_PATTERNS:
            if re.search(pattern, text_sample, re.IGNORECASE):
                return "royalty_statement"
        
        for pattern in self.CONTRACT_PATTERNS:
            if re.search(pattern, text_sample, re.IGNORECASE):
                return "contract"
        
        # Fallback to extension
        ext = Path(file_name).suffix.lower()
        if ext in self.EXTENSION_MAP:
            return self.EXTENSION_MAP[ext]
        
        return "unknown"

    def _normalize_doc_type(self, detected_type: str) -> str:
        """Normalize document type string."""
        mapping = {
            "split_sheet": "split_sheet",
            "splitsheet": "split_sheet",
            "split": "split_sheet",
            "royalty_statement": "royalty_statement",
            "royalty": "royalty_statement",
            "dsp_report": "dsp_report",
            "contract": "contract",
            "license": "license_contract",
            "agreement": "license_contract",
            "pdf": "unknown",
            "csv": "unknown",
            "xlsx": "unknown",
            "txt": "unknown",
        }
        return mapping.get(detected_type.lower(), "unknown")
